"""
Office and tabular parsers: docx, xlsx, pptx, csv.

These formats already carry explicit structure, so unlike the PDF parser
there is nothing to infer — the work is mapping their object models onto the
shared document model so they render through the same GFM serializer.

Each parser imports its library lazily: a deployment that only ever sees PDFs
should not need python-docx installed.
"""

from __future__ import annotations

import csv as _csv
import io
from pathlib import Path

from ..model import (
    Blockquote,
    Cell,
    Document,
    Heading,
    ListBlock,
    ListItem,
    Paragraph,
    Span,
    Table,
    text_span,
)

_W_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


class MissingDependency(RuntimeError):
    """Raised when the library for a format is not installed."""


# --------------------------------------------------------------------------
# docx
# --------------------------------------------------------------------------


def _docx_runs(paragraph) -> list[Span]:
    spans: list[Span] = []
    for run in paragraph.runs:
        if not run.text:
            continue
        style = (run.style.name if run.style else "") or ""
        spans.append(
            Span(
                run.text,
                bold=bool(run.bold),
                italic=bool(run.italic),
                code="Code" in style or "Mono" in style,
            )
        )
    if not spans and paragraph.text:
        spans = text_span(paragraph.text)
    return spans


def _docx_list_info(paragraph) -> tuple[bool, int] | None:
    """Return (ordered, indent_level) when the paragraph is a list item."""
    style = (paragraph.style.name if paragraph.style else "") or ""
    num_pr = paragraph._p.find(f"{_W_NS}pPr/{_W_NS}numPr")
    if num_pr is None and "List" not in style:
        return None
    level = 0
    if num_pr is not None:
        ilvl = num_pr.find(f"{_W_NS}ilvl")
        if ilvl is not None:
            level = int(ilvl.get(f"{_W_NS}val", "0"))
    ordered = "Number" in style
    return ordered, level


def _docx_heading_level(paragraph) -> int | None:
    style = (paragraph.style.name if paragraph.style else "") or ""
    if style == "Title":
        return 1
    if style.startswith("Heading"):
        tail = style.replace("Heading", "").strip()
        return int(tail) if tail.isdigit() else 2
    return None


def _docx_table(table) -> Table:
    rows: list[list[Cell]] = []
    for row in table.rows:
        cells: list[Cell] = []
        seen: set[int] = set()
        for cell in row.cells:
            # python-docx repeats the same object across a horizontal merge.
            ident = id(cell._tc)
            if ident in seen:
                if cells:
                    cells[-1].colspan += 1
                continue
            seen.add(ident)
            text = "\n".join(p.text for p in cell.paragraphs).strip()
            cells.append(Cell(spans=text_span(text) if text else []))
        rows.append(cells)
    if not rows:
        return Table()
    return Table(header=rows[0], rows=rows[1:])


def parse_docx(path: str | Path) -> Document:
    try:
        import docx  # type: ignore
        from docx.table import Table as DocxTable  # type: ignore
        from docx.text.paragraph import Paragraph as DocxParagraph  # type: ignore
        from docx.oxml.ns import qn  # type: ignore
    except ImportError as exc:  # pragma: no cover - environment dependent
        raise MissingDependency("python-docx is required for .docx") from exc

    document = docx.Document(str(path))
    doc = Document(meta={"format": "docx"})

    pending: list[tuple[bool, int, list[Span]]] = []

    def flush_list() -> None:
        """Emit buffered list paragraphs, nesting by indent level."""
        if not pending:
            return
        ordered = pending[0][0]
        root = ListBlock(ordered=ordered)
        stack: list[tuple[int, ListBlock]] = [(pending[0][1], root)]
        for is_ordered, level, spans in pending:
            while len(stack) > 1 and level < stack[-1][0]:
                stack.pop()
            if level > stack[-1][0]:
                child = ListBlock(ordered=is_ordered)
                parent_items = stack[-1][1].items
                if not parent_items:
                    parent_items.append(ListItem(blocks=[]))
                parent_items[-1].blocks.append(child)
                stack.append((level, child))
            stack[-1][1].items.append(ListItem(blocks=[Paragraph(spans)]))
        doc.add(root)
        pending.clear()

    for child in document.element.body.iterchildren():
        if child.tag == qn("w:p"):
            paragraph = DocxParagraph(child, document)
            spans = _docx_runs(paragraph)
            if not spans:
                continue
            list_info = _docx_list_info(paragraph)
            if list_info is not None:
                ordered, level = list_info
                if pending and pending[0][0] != ordered:
                    flush_list()
                pending.append((ordered, level, spans))
                continue
            flush_list()
            level = _docx_heading_level(paragraph)
            if level is not None:
                doc.add(Heading(level=level, spans=spans))
            else:
                doc.add(Paragraph(spans))
        elif child.tag == qn("w:tbl"):
            flush_list()
            doc.add(_docx_table(DocxTable(child, document)))

    flush_list()
    return doc


# --------------------------------------------------------------------------
# xlsx
# --------------------------------------------------------------------------


def _fmt_cell(value) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def parse_xlsx(path: str | Path, max_rows: int = 5000) -> Document:
    try:
        import openpyxl  # type: ignore
    except ImportError as exc:  # pragma: no cover - environment dependent
        raise MissingDependency("openpyxl is required for .xlsx") from exc

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    doc = Document(meta={"format": "xlsx"})
    try:
        for sheet in workbook.worksheets:
            doc.add(Heading(level=2, spans=text_span(sheet.title)))
            rows: list[list[Cell]] = []
            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                if i >= max_rows:
                    doc.warnings.append(
                        f"sheet '{sheet.title}' truncated at {max_rows} rows"
                    )
                    break
                values = [_fmt_cell(v) for v in row]
                while values and not values[-1]:
                    values.pop()
                if not values:
                    continue
                rows.append([Cell(spans=text_span(v) if v else []) for v in values])
            if not rows:
                doc.add(Paragraph(text_span("_(empty sheet)_")))
                continue
            doc.add(Table(header=rows[0], rows=rows[1:]))
    finally:
        workbook.close()
    return doc


# --------------------------------------------------------------------------
# pptx
# --------------------------------------------------------------------------


def parse_pptx(path: str | Path) -> Document:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except ImportError as exc:  # pragma: no cover - environment dependent
        raise MissingDependency("python-pptx is required for .pptx") from exc

    prs = Presentation(str(path))
    doc = Document(meta={"format": "pptx"})

    for index, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()
        doc.add(Heading(level=2, spans=text_span(title or f"Slide {index}")))

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_table:
                rows = [
                    [Cell(spans=text_span(c.text.strip()) if c.text.strip() else [])
                     for c in row.cells]
                    for row in shape.table.rows
                ]
                if rows:
                    doc.add(Table(header=rows[0], rows=rows[1:]))
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                doc.add(Paragraph(text_span(f"[image: {shape.name}]")))
                continue
            if not shape.has_text_frame:
                continue

            items: list[tuple[int, str]] = []
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs).strip() or para.text.strip()
                if text:
                    items.append((para.level or 0, text))
            if not items:
                continue
            if len(items) == 1 and items[0][0] == 0:
                doc.add(Paragraph(text_span(items[0][1])))
                continue

            root = ListBlock(ordered=False)
            stack: list[tuple[int, ListBlock]] = [(0, root)]
            for level, text in items:
                while len(stack) > 1 and level < stack[-1][0]:
                    stack.pop()
                if level > stack[-1][0]:
                    child = ListBlock(ordered=False)
                    if not stack[-1][1].items:
                        stack[-1][1].items.append(ListItem(blocks=[]))
                    stack[-1][1].items[-1].blocks.append(child)
                    stack.append((level, child))
                stack[-1][1].items.append(
                    ListItem(blocks=[Paragraph(text_span(text))])
                )
            doc.add(root)

        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                doc.add(Blockquote(blocks=[Paragraph(text_span(f"Notes: {notes}"))]))

    return doc


# --------------------------------------------------------------------------
# csv / plain text
# --------------------------------------------------------------------------


def parse_csv(path: str | Path, max_rows: int = 5000) -> Document:
    raw = Path(path).read_bytes()
    text = raw.decode("utf-8-sig", errors="replace")
    try:
        dialect = _csv.Sniffer().sniff(text[:4096], delimiters=",;\t|")
    except _csv.Error:
        dialect = _csv.excel

    doc = Document(meta={"format": "csv"})
    rows: list[list[Cell]] = []
    for i, row in enumerate(_csv.reader(io.StringIO(text), dialect)):
        if i >= max_rows:
            doc.warnings.append(f"truncated at {max_rows} rows")
            break
        rows.append([Cell(spans=text_span(v.strip()) if v.strip() else []) for v in row])
    if rows:
        doc.add(Table(header=rows[0], rows=rows[1:]))
    return doc


def parse_txt(path: str | Path) -> Document:
    text = Path(path).read_bytes().decode("utf-8-sig", errors="replace")
    doc = Document(meta={"format": "txt"})
    for chunk in text.split("\n\n"):
        chunk = chunk.strip()
        if chunk:
            doc.add(Paragraph(text_span(" ".join(chunk.split()))))
    return doc
